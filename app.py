from flask import Flask, request, jsonify, send_file
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pptx
import json
import uuid
import re
import os
from datetime import datetime
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image

book_map = {
    # 구약
    "창세기": "창", "창": "창", "창세": "창",
    "출애굽기": "출", "출": "출", "출애": "출",
    "레위기": "레", "레": "레", "레위": "레",
    "민수기": "민", "민": "민", "민수": "민",
    "신명기": "신", "신": "신", "신명": "신",
    "여호수아": "수", "수": "수", "여호": "수",
    "사사기": "삿", "삿": "삿", "판관": "삿",
    "룻기": "룻", "룻": "룻",
    "사무엘상": "삼상", "삼상": "삼상", "1사무": "삼상",
    "사무엘하": "삼하", "삼하": "삼하", "2사무": "삼하",
    "열왕기상": "왕상", "왕상": "왕상", "1열왕": "왕상",
    "열왕기하": "왕하", "왕하": "왕하", "2열왕": "왕하",
    "역대상": "대상", "대상": "대상", "1역대": "대상",
    "역대하": "대하", "대하": "대하", "2역대": "대하",
    "에스라": "스", "스": "스", "에즈": "스",
    "느헤미야": "느", "느": "느", "느헤": "느",
    "에스더": "에", "에": "에", "에스": "에",
    "욥기": "욥", "욥": "욥",
    "시편": "시", "시": "시", "시편": "시",
    "잠언": "잠", "잠": "잠", "잠언": "잠",
    "전도서": "전", "전": "전", "전도": "전",
    "아가": "아", "아": "아", "아가": "아",
    "이사야": "사", "사": "사", "이사": "사",
    "예레미야": "렘", "렘": "렘", "예레": "렘",
    "예레미야애가": "애", "애": "애", "애가": "애",
    "에스겔": "겔", "겔": "겔", "에제": "겔",
    "다니엘": "단", "단": "단", "다니": "단",
    "호세아": "호", "호": "호", "호세": "호",
    "요엘": "욜", "욜": "욜", "요엘": "욜",
    "아모스": "암", "암": "암", "아모": "암",
    "오바댜": "옵", "옵": "옵", "오바": "옵",
    "요나": "욘", "욘": "욘", "요나": "욘",
    "미가": "미", "미": "미", "미가": "미",
    "나훔": "나", "나": "나", "나훔": "나",
    "하박국": "합", "합": "합", "하바": "합",
    "스바냐": "습", "습": "습", "스바": "습",
    "학개": "학", "학": "학", "하깨": "학",
    "스가랴": "슥", "슥": "슥", "즈가": "슥",
    "말라기": "말", "말": "말", "말라": "말",

    # 신약
    "마태복음": "마", "마": "마", "마태": "마",
    "마가복음": "막", "막": "막", "마르": "막",
    "누가복음": "눅", "눅": "눅", "루가": "눅",
    "요한복음": "요", "요": "요", "요한": "요",
    "사도행전": "행", "행": "행", "사도": "행",
    "로마서": "롬", "롬": "롬", "로마": "롬",
    "고린도전서": "고전", "고전": "고전", "1고린": "고전",
    "고린도후서": "고후", "고후": "고후", "2고린": "고후",
    "갈라디아서": "갈", "갈": "갈", "갈라": "갈",
    "에베소서": "엡", "엡": "엡", "에페": "엡",
    "빌립보서": "빌", "빌": "빌", "필립": "빌",
    "골로새서": "골", "골": "골", "골로": "골",
    "데살로니가전서": "살전", "살전": "살전", "1데살": "살전",
    "데살로니가후서": "살후", "살후": "살후", "2데살": "살후",
    "디모데전서": "딤전", "딤전": "딤전", "1디모": "딤전",
    "디모데후서": "딤후", "딤후": "딤후", "2디모": "딤후",
    "디도서": "딛", "딛": "딛", "디도": "딛",
    "빌레몬서": "몬", "몬": "몬", "필레": "몬",
    "히브리서": "히", "히": "히", "히브": "히",
    "야고보서": "약", "약": "약", "야고": "약",
    "베드로전서": "벧전", "벧전": "벧전", "1베드": "벧전",
    "베드로후서": "벧후", "벧후": "벧후", "2베드": "벧후",
    "요한1서": "요일", "요일": "요일", "요1": "요일", "1요한": "요일",
    "요한2서": "요이", "요이": "요이", "요2": "요이", "2요한": "요이",
    "요한3서": "요삼", "요삼": "요삼", "요3": "요삼", "3요한": "요삼",
    "유다서": "유", "유": "유", "유다": "유",
    "요한계시록": "계", "계": "계", "묵시": "계"
}

short_to_full = {
    "창": "창세기",
    "출": "출애굽기",
    "레": "레위기",
    "민": "민수기",
    "신": "신명기",
    "수": "여호수아",
    "삿": "사사기",
    "룻": "룻기",
    "삼상": "사무엘상",
    "삼하": "사무엘하",
    "왕상": "열왕기상",
    "왕하": "열왕기하",
    "대상": "역대상",
    "대하": "역대하",
    "스": "에스라",
    "느": "느헤미야",
    "에": "에스더",
    "욥": "욥기",
    "시": "시편",
    "잠": "잠언",
    "전": "전도서",
    "아": "아가",
    "사": "이사야",
    "렘": "예레미야",
    "애": "예레미야애가",
    "겔": "에스겔",
    "단": "다니엘",
    "호": "호세아",
    "욜": "요엘",
    "암": "아모스",
    "옵": "오바댜",
    "욘": "요나",
    "미": "미가",
    "나": "나훔",
    "합": "하박국",
    "습": "스바냐",
    "학": "학개",
    "슥": "스가랴",
    "말": "말라기",

    "마": "마태복음",
    "막": "마가복음",
    "눅": "누가복음",
    "요": "요한복음",
    "행": "사도행전",
    "롬": "로마서",
    "고전": "고린도전서",
    "고후": "고린도후서",
    "갈": "갈라디아서",
    "엡": "에베소서",
    "빌": "빌립보서",
    "골": "골로새서",
    "살전": "데살로니가전서",
    "살후": "데살로니가후서",
    "딤전": "디모데전서",
    "딤후": "디모데후서",
    "딛": "디도서",
    "몬": "빌레몬서",
    "히": "히브리서",
    "약": "야고보서",
    "벧전": "베드로전서",
    "벧후": "베드로후서",
    "요일": "요한1서",
    "요이": "요한2서",
    "요삼": "요한3서",
    "유": "유다서",
    "계": "요한계시록"
}

app = Flask(__name__)

with open('bible.json', 'r', encoding='utf-8') as f:
    bible_data = json.load(f)

import re

def parse_verse_input(text):
    text = text.strip()
    text = text.replace(" ", "")
    text = text.replace("~", "-")
    text = text.rstrip(":")  # 겔45: 처럼 끝에 콜론 있어도 처리

    # 범위 입력: 겔44:15-45:12 / 겔44:15-20 / 겔44:15
    range_match = re.search(r'(\d+):(\d+)(?:-(?:(\d+):)?(\d+))?$', text)

    if range_match:
        book_input = text[:range_match.start()]
        start_chapter = int(range_match.group(1))
        start_verse = int(range_match.group(2))

        if range_match.group(4):
            end_chapter = int(range_match.group(3)) if range_match.group(3) else start_chapter
            end_verse = int(range_match.group(4))
        else:
            end_chapter = start_chapter
            end_verse = start_verse

    else:
        # 장 전체 입력: 겔45 / 에스겔45
        whole_chapter_match = re.search(r'(\d+)$', text)

        if not whole_chapter_match:
            raise ValueError("입력 형식 오류")

        book_input = text[:whole_chapter_match.start()]
        start_chapter = int(whole_chapter_match.group(1))
        start_verse = 1
        end_chapter = start_chapter
        end_verse = None

    book = book_map.get(book_input)

    if not book:
        raise ValueError(f"지원하지 않는 책 이름: {book_input}")

    return book, start_chapter, start_verse, end_chapter, end_verse


@app.route('/')
def home():
    return send_file('static/index.html')

@app.route('/preview', methods=['POST'])
def preview():
    data = request.json
    verse = data['verse']
    version = data['version']

    verses = get_bible_verses(verse, version)

    return jsonify(verses)

def hex_to_rgb_tuple(hex_color):
    hex_color = hex_color.replace("#", "")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_gradient_background(filename, start_hex="FFFFFF", end_hex="FFF7F3", width=1920, height=1080):
    start = hex_to_rgb_tuple(start_hex)
    end = hex_to_rgb_tuple(end_hex)

    img = Image.new("RGB", (width, height), start)

    for y in range(height):
        ratio = y / height
        r = int(start[0] * (1 - ratio) + end[0] * ratio)
        g = int(start[1] * (1 - ratio) + end[1] * ratio)
        b = int(start[2] * (1 - ratio) + end[2] * ratio)

        for x in range(width):
            img.putpixel((x, y), (r, g, b))

    img.save(filename)

@app.route('/generate', methods=['POST'])
def generate_ppt():
    try:
        data = request.json

        verse = data['verse']
        version = data['version']
        font_size = int(data['fontSize'])
        is_bold = data.get('isBold', True)

        verses = get_bible_verses(verse, version)

        prs = pptx.Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        ppt_font_name = "Malgun Gothic"

        text_color = RGBColor.from_string("FFFFFF")
        ref_box_color = RGBColor.from_string("54453D")

        gradient_file = f"gradient_{uuid.uuid4().hex}.png"
        create_gradient_background(
        gradient_file,
        start_hex="1C1616",
        end_hex="352115"
        )

        for v in verses:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 배경 그라데이션 이미지
            slide.shapes.add_picture(
                gradient_file,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height
            )

            # 절 이름 박스
            ref_box_width = Inches(3.8)
            ref_box_height = Inches(0.72)
            ref_box_x = (prs.slide_width - ref_box_width) / 2
            ref_box_y = Inches(0.45)

            ref_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                ref_box_x,
                ref_box_y,
                ref_box_width,
                ref_box_height
            )

            ref_shape.fill.solid()
            ref_shape.fill.fore_color.rgb = ref_box_color
            ref_shape.line.fill.background()

            ref_frame = ref_shape.text_frame
            ref_frame.clear()
            ref_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            ref_frame.margin_left = 0
            ref_frame.margin_right = 0
            ref_frame.margin_top = 0
            ref_frame.margin_bottom = 0

            ref_p = ref_frame.paragraphs[0]
            ref_p.alignment = PP_ALIGN.CENTER

            ref_run = ref_p.add_run()
            ref_run.text = v['reference']
            ref_run.font.name = ppt_font_name
            ref_run.font.size = Pt(24)
            ref_run.font.bold = True
            ref_run.font.color.rgb = text_color

            # 본문 텍스트
            text_box = slide.shapes.add_textbox(
             Inches(0.35),   # 왼쪽 여백 줄임
             Inches(1.75),
             Inches(12.25),  # 폭 넓힘
             Inches(4.5)
            )

            text_frame = text_box.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0

            p = text_frame.paragraphs[0]
            p.line_spacing = 1.15
            p.alignment = PP_ALIGN.CENTER

            run = p.add_run()
            run.text = v['text']
            run.font.name = ppt_font_name
            run.font.size = Pt(font_size)
            run.font.bold = is_bold
            run.font.color.rgb = text_color

        today = datetime.now().strftime("%Y-%m-%d")

        book, start_chapter, start_verse, end_chapter, end_verse = parse_verse_input(verse)
        full_book = short_to_full.get(book, book)

        if end_verse is None:
            filename_verse_part = f"{start_chapter}장전체"
        elif start_chapter == end_chapter:
            if start_verse == end_verse:
                filename_verse_part = f"{start_chapter}장{start_verse}절"
            else:
                filename_verse_part = f"{start_chapter}장{start_verse}-{end_verse}절"
        else:
            filename_verse_part = f"{start_chapter}장{start_verse}절-{end_chapter}장{end_verse}절"

        timestamp = datetime.now().strftime("%H%M%S")
        filename = f"{today}_{full_book}{filename_verse_part}_{timestamp}.pptx"
        prs.save(filename)

        return send_file(
            filename,
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        return jsonify({"error": str(e)}), 400
    
def format_text_for_ppt(text, font_size):
    words = text.split()
    lines = []
    current_line = ""

    # 폰트 크기에 따라 한 줄 최대 글자 수를 대충 조정
    if font_size >= 48:
        max_chars = 30
    elif font_size >= 42:
        max_chars = 34
    elif font_size >= 36:
        max_chars = 40
    elif font_size >= 30:
        max_chars = 48
    elif font_size >= 24:
        max_chars = 56
    else:
        max_chars = 64

    for word in words:
        test_line = current_line + (" " if current_line else "") + word

        if len(test_line) <= max_chars:
            current_line = test_line
        else:
            if current_line:
                lines.append(current_line)
            current_line = word

    if current_line:
        lines.append(current_line)

    return "\n".join(lines)

def get_bible_verses(verse_input, version):
    book, start_chapter, start_verse, end_chapter, end_verse = parse_verse_input(verse_input)

    if version not in bible_data:
        raise ValueError(f"지원하지 않는 성경 버전: {version}")

    verses = []
    full_book = short_to_full.get(book, book)

    # 해당 책/장에 존재하는 절 번호들 찾기
    def get_existing_verses(chapter):
        verse_nums = []

        prefix = f"{book} {chapter}:"

        for key in bible_data[version].keys():
            if key.startswith(prefix):
                try:
                    verse_num = int(key.split(":")[1])
                    verse_nums.append(verse_num)
                except:
                    pass

        return sorted(verse_nums)

    # 장 전체 입력인 경우: end_verse가 None
    if end_verse is None:
        existing = get_existing_verses(start_chapter)

        if not existing:
            raise ValueError(f"{full_book} {start_chapter}장을 찾을 수 없습니다")

        end_verse = existing[-1]

    for chapter in range(start_chapter, end_chapter + 1):
        existing = get_existing_verses(chapter)

        if not existing:
            continue

        chapter_start = start_verse if chapter == start_chapter else existing[0]
        chapter_end = end_verse if chapter == end_chapter else existing[-1]

        for v in range(chapter_start, chapter_end + 1):
            key = f"{book} {chapter}:{v}"
            text = bible_data[version].get(key)

            if text:
                reference = f"{full_book} {chapter}장 {v}절"

                verses.append({
                    "reference": reference,
                    "text": text
                })

    if not verses:
        raise ValueError("해당 구절을 찾을 수 없습니다")

    return verses

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5512)
